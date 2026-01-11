from pathlib import Path
from typing import Optional, Iterable, List, Any, Dict
from logger import GLOBAL_LOGGER as logger
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from utils.model_loader import ModelLoader
from utils.qdrant_vector_db import QdrantVDB
from utils.s3_operations import S3ReadUpload
from utils.config_loader import load_config

SUPPORTED_FILE_TYPES = [".txt", ".pdf", ".docx", ".md", ".pptx"]

class DataIngestion:

    def __init__(self):
        ### Load Embedding and LLM models:
        model_load = ModelLoader()
        self.embeddings = model_load.load_embeddings()
        
        ### Initialize Qdrant Vector DB
        self.vector_db = QdrantVDB()

        ### Initialize S3 Operations
        self.s3_ops = S3ReadUpload()

        ### Load S3 configuration
        config = load_config("config.yaml")
        s3_config = config.get('AWS-S3', {})
        self.bucket_name = s3_config.get('bucket_name')
        self.object_prefix = s3_config.get('preindex_folder_name')

    def ingest_files(self, file_paths: List[Path], user_name: str) -> None:
        documents = []
        for file_path in file_paths:
            if file_path.suffix.lower() in SUPPORTED_FILE_TYPES:
                content = self._read_file_content(file_path)
                doc = Document(page_content=content, metadata={"source": str(file_path)})
                documents.append(doc)
                object_name = f"{self.object_prefix}/{user_name}/{file_path.name}"
                self.s3_ops.upload_file_to_s3(file_name=str(file_path), bucket_name=self.bucket_name, object_name=object_name)
                logger.info(f"File uploaded to S3: {object_name}")
            else:
                print(f"Unsupported file type: {file_path.suffix} for file {file_path}")

        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            separators=["\n\n", "\n", " ", ""]
        )
        split_docs = text_splitter.split_documents(documents)

        self.vector_db.create_vector_store(
            embedding=self.embeddings,
            collection_name=user_name,
            documents=split_docs
        )
        print(f"Ingested {len(split_docs)} documents into collection '{user_name}'.")

    def _read_file_content(self, file_path: Path) -> str:
        if file_path.suffix.lower() == ".txt":
            content = self._read_txt(file_path)
        elif file_path.suffix.lower() == ".pdf":
            content = self._read_pdf(file_path)
        elif file_path.suffix.lower() == ".docx":
            content = self._read_docx(file_path)
        elif file_path.suffix.lower() == ".md":
            content = self._read_md(file_path)
        elif file_path.suffix.lower() == ".pptx":
            content = self._read_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_path.suffix}")
        return content
    
    def _read_txt(self, file_path: Path) -> str:
        with open(file_path, "r", encoding="utf-8") as f:
            logger.info(f"File read (txt): {file_path}")
            return f.read()

    def _read_pdf(self, file_path: Path) -> str:
        import fitz  # PyMuPDF
        text_chunks = []
        with fitz.open(file_path) as doc:
            for page_num in range(doc.page_count):
                page = doc.load_page(page_num)
                text_chunks.append(f"\n--- Page {page_num + 1} ---\n{page.get_text()}")
        
        return "\n".join(text_chunks)
    
    def _read_docx(self, file_path: Path) -> str:
        import docx2txt as docx
        text = docx.process(file_path)
        logger.info(f"File read (docx): {file_path}")
        return text
    
    def _read_md(self, file_path: Path) -> str:
        with open(file_path, "r", encoding="utf-8") as f:
            logger.info(f"File read (md): {file_path}")
            return f.read()
        
    def _read_pptx(self, file_path: Path) -> str:
        from pptx import Presentation
        prs = Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        logger.info(f"File read (pptx): {file_path}")
        return "\n".join(full_text)

if __name__ == "__main__":
    data_ingestion = DataIngestion()
    test_files = [Path("./data/text.txt"), Path("./data/Arindam_dec_2025.docx"), Path("./data/Rudy-2025.pdf"), Path("./data/TJX Case Study Jan 2026.pptx"), Path("./data/SETUP.md")]
    data_ingestion.ingest_files(file_paths=test_files, user_name="Arindam")